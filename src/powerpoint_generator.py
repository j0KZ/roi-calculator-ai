"""
Professional PowerPoint Generator for ROI Calculator
Creates comprehensive PowerPoint presentations with multiple templates and data visualizations
"""

import os
import json
import tempfile
from datetime import datetime, timedelta
from pathlib import Path
import logging
from typing import Dict, Any, List, Optional, Tuple
import base64
from io import BytesIO

# PowerPoint generation imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# Chart generation
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import seaborn as sns
import numpy as np
from PIL import Image as PILImage

logger = logging.getLogger(__name__)

class PowerPointGenerator:
    """Professional PowerPoint generator with multiple template support"""
    
    # Template configurations
    TEMPLATES = {
        'executive': {
            'name': 'Executive',
            'description': 'Clean, minimal design for C-level presentations',
            'primary_color': RGBColor(0x1B, 0x36, 0x5D),
            'secondary_color': RGBColor(0x5D, 0x73, 0x7E),
            'accent_color': RGBColor(0x00, 0xB8, 0x94),
            'font_name': 'Calibri',
            'title_size': Pt(28),
            'subtitle_size': Pt(18),
            'body_size': Pt(14)
        },
        'sales': {
            'name': 'Sales',
            'description': 'Dynamic, persuasive design for sales presentations',
            'primary_color': RGBColor(0xFF, 0x6B, 0x6B),
            'secondary_color': RGBColor(0x4E, 0xCD, 0xC4),
            'accent_color': RGBColor(0xFF, 0xE6, 0x6D),
            'font_name': 'Segoe UI',
            'title_size': Pt(32),
            'subtitle_size': Pt(20),
            'body_size': Pt(16)
        },
        'technical': {
            'name': 'Technical',
            'description': 'Professional design for detailed technical analysis',
            'primary_color': RGBColor(0x6C, 0x5C, 0xE7),
            'secondary_color': RGBColor(0xA2, 0x9B, 0xFE),
            'accent_color': RGBColor(0x00, 0xB8, 0x94),
            'font_name': 'Arial',
            'title_size': Pt(24),
            'subtitle_size': Pt(16),
            'body_size': Pt(12)
        }
    }
    
    def __init__(self, roi_results: Dict[str, Any], company_config: Optional[Dict] = None):
        """
        Initialize PowerPoint generator
        
        Args:
            roi_results: ROI calculation results
            company_config: Company branding configuration
        """
        self.roi_results = roi_results
        self.company_config = company_config or self._get_default_company_config()
        self.generated_charts = {}
        
        # Set up matplotlib for professional charts
        plt.style.use('seaborn-v0_8-whitegrid')
        sns.set_palette("husl")
        
        # Create temp directory for charts
        self.temp_dir = tempfile.mkdtemp(prefix='roi_powerpoint_')
        
    def _get_default_company_config(self) -> Dict:
        """Get default company configuration"""
        return {
            'name': 'ROI Solutions Inc.',
            'address': '123 Business Ave, Suite 100\nBusiness City, BC 12345',
            'phone': '+1 (555) 123-4567',
            'email': 'contact@roisolutions.com',
            'website': 'www.roisolutions.com',
            'logo_path': None,
            'tagline': 'Maximizing Your Investment Returns'
        }
    
    def generate_presentation(self, template_name: str = 'executive', 
                            custom_config: Optional[Dict] = None) -> str:
        """
        Generate PowerPoint presentation
        
        Args:
            template_name: Template to use ('executive', 'sales', 'technical')
            custom_config: Custom configuration overrides
            
        Returns:
            File path to generated presentation
        """
        try:
            # Apply custom configuration
            if custom_config:
                self.company_config.update(custom_config)
            
            # Create presentation
            prs = Presentation()
            
            # Set slide size (16:9 widescreen)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Get template config
            template_config = self.TEMPLATES.get(template_name, self.TEMPLATES['executive'])
            
            # Generate all charts first
            self._generate_all_charts(template_config)
            
            # Build slides
            self._add_title_slide(prs, template_config)
            self._add_executive_summary_slide(prs, template_config)
            self._add_current_state_slide(prs, template_config)
            self._add_roi_projections_slide(prs, template_config)
            self._add_cost_breakdown_slide(prs, template_config)
            self._add_savings_timeline_slide(prs, template_config)
            self._add_implementation_timeline_slide(prs, template_config)
            self._add_risk_assessment_slide(prs, template_config)
            self._add_next_steps_slide(prs, template_config)
            
            # Save presentation
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            client_name = self.roi_results.get('inputs', {}).get('company_name', 'Client')
            safe_client = ''.join(c for c in client_name if c.isalnum() or c in (' ', '-', '_')).strip()
            
            filename = f'ROI_Presentation_{safe_client}_{template_name}_{timestamp}.pptx'
            filepath = os.path.join(self.temp_dir, filename)
            
            prs.save(filepath)
            
            logger.info(f"Generated PowerPoint presentation: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint presentation: {str(e)}")
            raise
    
    def _generate_all_charts(self, template_config: Dict) -> None:
        """Generate all charts for the presentation"""
        try:
            self.generated_charts = {
                'roi_timeline': self._create_roi_timeline_chart(template_config),
                'cost_breakdown': self._create_cost_breakdown_chart(template_config),
                'savings_projection': self._create_savings_projection_chart(template_config),
                'investment_returns': self._create_investment_returns_chart(template_config),
                'risk_analysis': self._create_risk_analysis_chart(template_config),
                'implementation_gantt': self._create_implementation_gantt_chart(template_config)
            }
        except Exception as e:
            logger.error(f"Error generating charts: {str(e)}")
            self.generated_charts = {}
    
    def _create_roi_timeline_chart(self, template_config: Dict) -> str:
        """Create ROI timeline visualization"""
        fig, ax = plt.subplots(figsize=(12, 6))
        
        # Extract projection data
        projections = self.roi_results.get('projections', {})
        years = []
        roi_values = []
        cumulative_savings = []
        
        for year in ['year_1', 'year_2', 'year_3']:
            if year in projections:
                years.append(int(year.split('_')[1]))
                roi_values.append(projections[year].get('roi_percentage', 0))
                cumulative_savings.append(projections[year].get('cumulative_savings', 0))
        
        if not years:
            # Use default data if projections not available
            years = [1, 2, 3]
            base_roi = self.roi_results.get('roi_metrics', {}).get('first_year_roi', 0.25)
            roi_values = [base_roi, base_roi * 1.15, base_roi * 1.25]
            base_savings = self.roi_results.get('roi_metrics', {}).get('annual_savings', 50000)
            cumulative_savings = [base_savings, base_savings * 2.1, base_savings * 3.3]
        
        # Convert RGB to matplotlib color
        primary_color = template_config['primary_color']
        secondary_color = template_config['secondary_color']
        
        # Extract RGB values from RGBColor tuple
        primary_hex = f"#{primary_color[0]:02x}{primary_color[1]:02x}{primary_color[2]:02x}"
        secondary_hex = f"#{secondary_color[0]:02x}{secondary_color[1]:02x}{secondary_color[2]:02x}"
        
        # Create dual-axis chart
        ax2 = ax.twinx()
        
        # ROI percentage line
        line1 = ax.plot(years, [r * 100 for r in roi_values], marker='o', linewidth=4, 
                       color=primary_hex, label='ROI %', markersize=8)
        ax.set_ylabel('ROI Percentage (%)', fontsize=14, fontweight='bold')
        ax.set_ylim(0, max([r * 100 for r in roi_values]) * 1.3)
        
        # Cumulative savings bars
        bars = ax2.bar(years, cumulative_savings, alpha=0.7, width=0.6,
                      color=secondary_hex, label='Cumulative Savings ($)')
        ax2.set_ylabel('Cumulative Savings ($)', fontsize=14, fontweight='bold')
        
        # Formatting
        ax.set_xlabel('Year', fontsize=14, fontweight='bold')
        ax.set_title('ROI Timeline & Cumulative Savings', fontsize=18, fontweight='bold', pad=20)
        ax.set_xticks(years)
        ax.grid(True, alpha=0.3)
        
        # Format currency labels
        ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
        
        # Add value labels on bars and line points
        for bar, value in zip(bars, cumulative_savings):
            height = bar.get_height()
            ax2.text(bar.get_x() + bar.get_width()/2., height + height*0.02,
                    f'${value:,.0f}', ha='center', va='bottom', fontweight='bold', fontsize=11)
        
        for year, roi in zip(years, roi_values):
            ax.text(year, roi * 100 + max([r * 100 for r in roi_values]) * 0.05,
                   f'{roi * 100:.1f}%', ha='center', va='bottom', fontweight='bold', fontsize=11)
        
        # Combined legend
        lines1, labels1 = ax.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax.legend(lines1 + lines2, labels1 + labels2, loc='upper left', frameon=True, 
                 fancybox=True, shadow=True, fontsize=12)
        
        plt.tight_layout()
        
        chart_path = os.path.join(self.temp_dir, 'roi_timeline_chart.png')
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return chart_path
    
    def _create_cost_breakdown_chart(self, template_config: Dict) -> str:
        """Create cost breakdown donut chart"""
        fig, ax = plt.subplots(figsize=(10, 8))
        
        # Extract cost data
        inputs = self.roi_results.get('inputs', {})
        costs = {
            'Labor': inputs.get('labor_costs', 0) * 12,
            'Shipping': inputs.get('shipping_costs', 0) * 12,
            'Errors': inputs.get('error_costs', 0) * 12,
            'Inventory': inputs.get('inventory_costs', 0) * 12
        }
        
        # Filter out zero costs
        costs = {k: v for k, v in costs.items() if v > 0}
        
        if not costs:
            costs = {'Sample Costs': 100000}
        
        # Create donut chart with professional colors
        colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DDA0DD']
        
        # Create donut (pie with hole)
        wedges, texts, autotexts = ax.pie(costs.values(), labels=costs.keys(), 
                                         autopct='%1.1f%%', startangle=90,
                                         colors=colors[:len(costs)], 
                                         explode=[0.03] * len(costs),
                                         wedgeprops=dict(width=0.7))
        
        # Style the text
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
            autotext.set_fontsize(12)
        
        for text in texts:
            text.set_fontsize(14)
            text.set_fontweight('bold')
        
        ax.set_title('Annual Cost Breakdown', fontsize=18, fontweight='bold', pad=20)
        
        # Add total cost in center
        total_cost = sum(costs.values())
        ax.text(0, 0, f'Total Annual\nOperational Costs\n${total_cost:,.0f}', 
               ha='center', va='center', fontsize=16, fontweight='bold',
               bbox=dict(boxstyle='round,pad=0.8', facecolor='white', alpha=0.9, edgecolor='gray'))
        
        plt.tight_layout()
        
        chart_path = os.path.join(self.temp_dir, 'cost_breakdown_chart.png')
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return chart_path
    
    def _create_savings_projection_chart(self, template_config: Dict) -> str:
        """Create 36-month savings projection chart"""
        fig, ax = plt.subplots(figsize=(14, 8))
        
        # Generate monthly data
        monthly_savings = self.roi_results.get('roi_metrics', {}).get('monthly_savings', 5000)
        months = list(range(1, 37))
        
        # Model progressive savings (efficiency improvements over time)
        base_savings = [monthly_savings * (1 + 0.03 * (m-1)/12) for m in months]  # 3% annual improvement
        
        # Add seasonal variance
        seasonal_factor = [1 + 0.1 * np.sin((m-1) * 2 * np.pi / 12) for m in months]
        actual_savings = [s * f for s, f in zip(base_savings, seasonal_factor)]
        
        # Calculate cumulative
        cumulative_savings = np.cumsum(actual_savings)
        
        # Convert template colors
        primary_color = template_config['primary_color']
        accent_color = template_config['accent_color']
        
        primary_hex = f"#{primary_color[0]:02x}{primary_color[1]:02x}{primary_color[2]:02x}"
        accent_hex = f"#{accent_color[0]:02x}{accent_color[1]:02x}{accent_color[2]:02x}"
        
        # Plot cumulative savings
        ax.plot(months, cumulative_savings, linewidth=4, color=primary_hex, 
               label='Cumulative Savings', marker='o', markersize=4, markevery=3)
        ax.fill_between(months, 0, cumulative_savings, alpha=0.3, color=primary_hex)
        
        # Add monthly savings as secondary y-axis
        ax2 = ax.twinx()
        ax2.bar(months, actual_savings, alpha=0.6, color=accent_hex, 
               width=0.8, label='Monthly Savings')
        
        # Formatting
        ax.set_xlabel('Month', fontsize=14, fontweight='bold')
        ax.set_ylabel('Cumulative Savings ($)', fontsize=14, fontweight='bold', color=primary_hex)
        ax2.set_ylabel('Monthly Savings ($)', fontsize=14, fontweight='bold', color=accent_hex)
        ax.set_title('36-Month Savings Projection', fontsize=18, fontweight='bold', pad=20)
        
        # Format axes
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
        ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
        
        # Add milestone annotations
        for year in [12, 24, 36]:
            if year <= len(cumulative_savings):
                value = cumulative_savings[year-1]
                ax.annotate(f'Year {year//12}\n${value:,.0f}', 
                           xy=(year, value), xytext=(10, 20),
                           textcoords='offset points', fontsize=12, fontweight='bold',
                           bbox=dict(boxstyle='round,pad=0.4', facecolor='yellow', alpha=0.8),
                           arrowprops=dict(arrowstyle='->', color='black', lw=2))
        
        ax.grid(True, alpha=0.3)
        ax.set_xlim(0, 37)
        
        # Legends
        ax.legend(loc='upper left', fontsize=12)
        ax2.legend(loc='upper right', fontsize=12)
        
        plt.tight_layout()
        
        chart_path = os.path.join(self.temp_dir, 'savings_projection_chart.png')
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return chart_path
    
    def _create_investment_returns_chart(self, template_config: Dict) -> str:
        """Create investment vs returns waterfall chart"""
        fig, ax = plt.subplots(figsize=(12, 8))
        
        investment = self.roi_results.get('inputs', {}).get('service_investment', 50000)
        annual_savings = self.roi_results.get('roi_metrics', {}).get('annual_savings', 60000)
        
        # Waterfall chart data
        categories = ['Initial\nInvestment', 'Year 1\nSavings', 'Year 2\nSavings', 'Year 3\nSavings', 'Net\nReturn']
        values = [-investment, annual_savings, annual_savings * 1.1, annual_savings * 1.2, 0]
        
        # Calculate cumulative for waterfall effect
        cumulative = [0]
        for i, val in enumerate(values[:-1]):
            cumulative.append(cumulative[-1] + val)
        
        # Net return is the final cumulative
        values[-1] = cumulative[-1]
        cumulative[-1] = 0  # Start net return from zero
        
        # Colors: red for investment, green for positive values
        colors = ['#FF4444', '#44AA44', '#44AA44', '#44AA44', '#2266AA']
        
        # Create waterfall bars
        bars = []
        for i, (cat, val, cum) in enumerate(zip(categories, values, cumulative)):
            if i == len(categories) - 1:  # Net return bar
                bar = ax.bar(cat, val, bottom=0, color=colors[i], alpha=0.8)
            else:
                bar = ax.bar(cat, abs(val), bottom=cum if val > 0 else cum + val, 
                           color=colors[i], alpha=0.8)
            bars.append(bar)
            
            # Add value labels
            label_y = cum + val/2 if i < len(categories) - 1 else val/2
            ax.text(i, label_y, f'${abs(val):,.0f}', ha='center', va='center',
                   fontweight='bold', fontsize=12, color='white')
        
        # Add connecting lines for waterfall effect
        for i in range(len(categories) - 2):
            start_y = cumulative[i+1]
            end_y = cumulative[i+1]
            ax.plot([i+0.4, i+1-0.4], [start_y, end_y], 'k--', alpha=0.5, linewidth=1)
        
        # Break-even line
        ax.axhline(y=0, color='black', linestyle='-', linewidth=2, alpha=0.8)
        
        # Formatting
        ax.set_ylabel('Amount ($)', fontsize=14, fontweight='bold')
        ax.set_title('Investment vs Returns Analysis (3-Year Waterfall)', fontsize=18, fontweight='bold', pad=20)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
        
        # Add ROI percentage
        roi_3_year = (values[-1] / investment) * 100 if investment > 0 else 0
        ax.text(0.02, 0.98, f'3-Year ROI: {roi_3_year:.1f}%', transform=ax.transAxes, 
               fontsize=16, fontweight='bold', va='top',
               bbox=dict(boxstyle='round,pad=0.5', facecolor='lightgreen', alpha=0.8))
        
        ax.grid(True, alpha=0.3, axis='y')
        plt.xticks(rotation=0, ha='center')
        plt.tight_layout()
        
        chart_path = os.path.join(self.temp_dir, 'investment_returns_chart.png')
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return chart_path
    
    def _create_risk_analysis_chart(self, template_config: Dict) -> str:
        """Create risk assessment radar chart"""
        fig, ax = plt.subplots(figsize=(10, 10), subplot_kw=dict(projection='polar'))
        
        # Risk categories and scores (1-5 scale, where 5 is lowest risk)
        categories = ['Market Risk', 'Technology Risk', 'Implementation Risk', 
                     'Financial Risk', 'Operational Risk', 'Competitive Risk']
        
        # Sample risk scores (in real implementation, these would be calculated)
        risk_scores = [4.2, 4.5, 3.8, 4.0, 4.3, 3.9]
        
        # Number of categories
        N = len(categories)
        
        # Compute angles for each category
        angles = [n / float(N) * 2 * np.pi for n in range(N)]
        angles += angles[:1]  # Complete the circle
        
        # Add scores
        risk_scores += risk_scores[:1]  # Complete the circle
        
        # Plot
        primary_color = template_config['primary_color']
        accent_color = template_config['accent_color']
        
        primary_hex = f"#{primary_color[0]:02x}{primary_color[1]:02x}{primary_color[2]:02x}"
        accent_hex = f"#{accent_color[0]:02x}{accent_color[1]:02x}{accent_color[2]:02x}"
        
        ax.plot(angles, risk_scores, 'o-', linewidth=3, color=primary_hex, markersize=8)
        ax.fill(angles, risk_scores, alpha=0.25, color=primary_hex)
        
        # Add category labels
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(categories, fontsize=12)
        
        # Set y-axis
        ax.set_ylim(0, 5)
        ax.set_yticks([1, 2, 3, 4, 5])
        ax.set_yticklabels(['High Risk', 'Med-High', 'Medium', 'Med-Low', 'Low Risk'], fontsize=10)
        ax.grid(True)
        
        # Add title
        ax.set_title('Risk Assessment Profile\n(Higher values = Lower risk)', 
                    fontsize=16, fontweight='bold', pad=30)
        
        # Add average risk score
        avg_risk = np.mean(risk_scores[:-1])
        risk_level = 'Low' if avg_risk >= 4.0 else 'Medium' if avg_risk >= 3.0 else 'High'
        color = 'green' if avg_risk >= 4.0 else 'orange' if avg_risk >= 3.0 else 'red'
        
        ax.text(0.5, 0.02, f'Overall Risk Level: {risk_level} ({avg_risk:.1f}/5.0)', 
               transform=ax.transAxes, ha='center', va='bottom', fontsize=14, fontweight='bold',
               bbox=dict(boxstyle='round,pad=0.5', facecolor=color, alpha=0.7, edgecolor='black'))
        
        plt.tight_layout()
        
        chart_path = os.path.join(self.temp_dir, 'risk_analysis_chart.png')
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return chart_path
    
    def _create_implementation_gantt_chart(self, template_config: Dict) -> str:
        """Create implementation timeline Gantt chart"""
        fig, ax = plt.subplots(figsize=(14, 8))
        
        # Implementation phases
        phases = [
            ('Project Kickoff', 1, 2, '#FF6B6B'),
            ('System Analysis', 2, 3, '#4ECDC4'),
            ('Configuration', 3, 6, '#45B7D1'),
            ('Integration', 5, 8, '#96CEB4'),
            ('Testing', 7, 10, '#FFEAA7'),
            ('Training', 8, 11, '#DDA0DD'),
            ('Go-Live', 10, 12, '#98FB98'),
            ('Optimization', 12, 24, '#F0E68C')
        ]
        
        # Create Gantt bars
        for i, (phase, start, end, color) in enumerate(phases):
            ax.barh(i, end - start, left=start, height=0.6, 
                   color=color, alpha=0.8, edgecolor='black', linewidth=1)
            
            # Add phase labels
            ax.text(start + (end - start) / 2, i, phase, 
                   ha='center', va='center', fontweight='bold', fontsize=10)
        
        # Formatting
        ax.set_xlabel('Weeks', fontsize=14, fontweight='bold')
        ax.set_ylabel('Implementation Phases', fontsize=14, fontweight='bold')
        ax.set_title('Implementation Timeline (Gantt Chart)', fontsize=18, fontweight='bold', pad=20)
        
        # Set y-axis
        ax.set_yticks(range(len(phases)))
        ax.set_yticklabels([phase[0] for phase in phases])
        ax.invert_yaxis()  # Phases from top to bottom
        
        # Set x-axis
        ax.set_xlim(0, 26)
        ax.set_xticks(range(0, 27, 2))
        
        # Add milestones
        milestones = [(4, 'System Ready'), (12, 'Go-Live'), (24, 'Full ROI')]
        for week, milestone in milestones:
            ax.axvline(x=week, color='red', linestyle='--', alpha=0.7, linewidth=2)
            ax.text(week, -0.8, milestone, ha='center', va='top', fontweight='bold', 
                   fontsize=10, rotation=45,
                   bbox=dict(boxstyle='round,pad=0.3', facecolor='yellow', alpha=0.8))
        
        ax.grid(True, alpha=0.3, axis='x')
        plt.tight_layout()
        
        chart_path = os.path.join(self.temp_dir, 'implementation_gantt_chart.png')
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return chart_path
    
    def _add_title_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add title slide with company info"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Title
        title.text = "ROI Analysis & Business Case"
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = template_config['title_size']
        title_paragraph.font.name = template_config['font_name']
        title_paragraph.font.color.rgb = template_config['primary_color']
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        client_name = self.roi_results.get('inputs', {}).get('company_name', 'Valued Client')
        subtitle_text = f"""Prepared for: {client_name}
        
Prepared by: {self.company_config['name']}
{self.company_config.get('tagline', 'Professional ROI Analysis')}

Date: {datetime.now().strftime('%B %d, %Y')}"""
        
        subtitle.text = subtitle_text
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = template_config['subtitle_size']
        subtitle_paragraph.font.name = template_config['font_name']
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add company contact info at bottom
        textbox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(1))
        text_frame = textbox.text_frame
        text_frame.text = f"{self.company_config['phone']} | {self.company_config['email']} | {self.company_config['website']}"
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.name = template_config['font_name']
        p.font.color.rgb = template_config['secondary_color']
        p.alignment = PP_ALIGN.CENTER
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"""
Welcome to the ROI analysis presentation for {client_name}. This presentation will demonstrate:
- Current operational challenges and costs
- Projected return on investment
- Implementation timeline and risk assessment
- Next steps for moving forward

Key talking points:
- Emphasize the strong ROI projections
- Address any concerns about implementation risk
- Highlight the payback period and long-term benefits
"""
    
    def _add_executive_summary_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add executive summary slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        self._format_title(title, template_config)
        
        # Key metrics
        investment = self.roi_results.get('inputs', {}).get('service_investment', 0)
        roi = self.roi_results.get('roi_metrics', {}).get('first_year_roi', 0) * 100
        annual_savings = self.roi_results.get('roi_metrics', {}).get('annual_savings', 0)
        payback_months = self.roi_results.get('roi_metrics', {}).get('payback_period_months', 0)
        
        # Add content
        content = slide.placeholders[1]
        content.text = f"""Investment Overview:
• Initial Investment: ${investment:,.2f}
• First Year ROI: {roi:.1f}%
• Annual Savings: ${annual_savings:,.2f}
• Payback Period: {payback_months:.1f} months

Key Benefits:
• Operational efficiency improvements
• Cost reduction across all major categories
• Enhanced productivity and reduced errors
• Scalable solution for future growth

Strategic Impact:
• Competitive advantage through operational excellence
• Improved cash flow and profitability
• Foundation for digital transformation initiatives"""
        
        self._format_content(content, template_config)
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = """
Executive Summary talking points:
- Start with the compelling ROI numbers
- Emphasize that this is more than just cost savings - it's about competitive advantage
- Highlight the relatively short payback period
- Connect operational improvements to strategic business goals
- Be prepared to discuss how these projections were calculated
"""
    
    def _add_current_state_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add current state analysis slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "Current State Analysis"
        self._format_title_textbox(title_frame, template_config)
        
        # Current costs breakdown
        inputs = self.roi_results.get('inputs', {})
        
        # Left side - text content
        left_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        left_frame = left_textbox.text_frame
        
        annual_labor = inputs.get('labor_costs', 0) * 12
        annual_shipping = inputs.get('shipping_costs', 0) * 12
        annual_errors = inputs.get('error_costs', 0) * 12
        annual_inventory = inputs.get('inventory_costs', 0) * 12
        total_annual = annual_labor + annual_shipping + annual_errors + annual_inventory
        
        left_frame.text = f"""Current Annual Operational Costs:

Labor Costs: ${annual_labor:,.2f}
• Manual processes and inefficiencies
• Error correction and rework
• Overtime and temporary staff

Shipping & Logistics: ${annual_shipping:,.2f}
• Suboptimal routing and carrier selection
• Individual shipments vs. bulk opportunities

Error Costs: ${annual_errors:,.2f}
• Returns processing and refunds
• Customer service overhead
• Inventory writeoffs

Inventory Management: ${annual_inventory:,.2f}
• Excess inventory carrying costs
• Stockouts and opportunity costs
• Manual inventory tracking

Total Annual Impact: ${total_annual:,.2f}"""
        
        self._format_content_textbox(left_frame, template_config, font_size=Pt(12))
        
        # Right side - cost breakdown chart
        if 'cost_breakdown' in self.generated_charts:
            chart_path = self.generated_charts['cost_breakdown']
            slide.shapes.add_picture(chart_path, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5))
        
        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = """
Current State Analysis talking points:
- Walk through each cost category with specific examples
- Relate costs to business impact (customer satisfaction, competitiveness)
- Emphasize that these are ongoing, recurring costs
- Ask client to validate these figures against their experience
- Set up the 'pain points' that our solution addresses
"""
    
    def _add_roi_projections_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add ROI projections slide with timeline chart"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "ROI Projections & Timeline"
        self._format_title_textbox(title_frame, template_config)
        
        # Add ROI timeline chart
        if 'roi_timeline' in self.generated_charts:
            chart_path = self.generated_charts['roi_timeline']
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        
        # Key metrics summary at bottom
        metrics_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.8))
        metrics_frame = metrics_box.text_frame
        
        roi_y3 = self.roi_results.get('projections', {}).get('year_3', {}).get('roi_percentage', 0) * 100
        npv = self.roi_results.get('financial_metrics', {}).get('npv', 0)
        irr = self.roi_results.get('financial_metrics', {}).get('irr', 0) * 100
        
        metrics_frame.text = f"3-Year ROI: {roi_y3:.1f}% | NPV: ${npv:,.2f} | IRR: {irr:.1f}% | Low Risk Profile"
        
        p = metrics_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.name = template_config['font_name']
        p.font.bold = True
        p.font.color.rgb = template_config['accent_color']
        p.alignment = PP_ALIGN.CENTER
        
        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = """
ROI Projections talking points:
- Walk through the timeline showing when benefits begin
- Explain the cumulative nature of savings over time
- Compare to alternative investment opportunities (NPV/IRR)
- Address any questions about the growth assumptions
- Emphasize conservative projections with upside potential
"""
    
    def _add_cost_breakdown_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add detailed cost breakdown slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "Cost Breakdown & Savings Opportunities"
        self._format_title_textbox(title_frame, template_config)
        
        # Split into two columns
        # Left column - detailed breakdown
        left_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5))
        left_frame = left_textbox.text_frame
        
        savings = self.roi_results.get('savings', {})
        
        left_frame.text = f"""Savings by Category:

Labor Optimization ({savings.get('labor', {}).get('percentage', 0)*100:.0f}% reduction)
• Process automation and workflow optimization
• Reduced error correction time
• Annual Savings: ${savings.get('labor', {}).get('annual', 0):,.2f}

Shipping Optimization ({savings.get('shipping', {}).get('percentage', 0)*100:.0f}% reduction)
• Carrier rate optimization
• Route planning and consolidation
• Annual Savings: ${savings.get('shipping', {}).get('annual', 0):,.2f}

Error Elimination ({savings.get('errors', {}).get('percentage', 0)*100:.0f}% reduction)
• Quality control systems
• Automated validation processes
• Annual Savings: ${savings.get('errors', {}).get('annual', 0):,.2f}

Inventory Optimization ({savings.get('inventory', {}).get('percentage', 0)*100:.0f}% reduction)
• Demand forecasting improvements
• Just-in-time inventory management
• Annual Savings: ${savings.get('inventory', {}).get('annual', 0):,.2f}

Total Annual Savings: ${savings.get('annual_total', 0):,.2f}"""
        
        self._format_content_textbox(left_frame, template_config, font_size=Pt(11))
        
        # Right column - donut chart
        if 'cost_breakdown' in self.generated_charts:
            chart_path = self.generated_charts['cost_breakdown']
            slide.shapes.add_picture(chart_path, Inches(7), Inches(1.5), Inches(5.5), Inches(5))
        
        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = """
Cost Breakdown talking points:
- Explain how each savings category is achieved
- Provide specific examples relevant to client's industry
- Address any skepticism about savings percentages
- Connect savings to operational improvements
- Discuss implementation approach for each category
"""
    
    def _add_savings_timeline_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add 36-month savings timeline slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "36-Month Savings Projection"
        self._format_title_textbox(title_frame, template_config)
        
        # Add savings projection chart
        if 'savings_projection' in self.generated_charts:
            chart_path = self.generated_charts['savings_projection']
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        
        # Key insights at bottom
        insights_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
        insights_frame = insights_box.text_frame
        
        insights_frame.text = "Key Insights: Progressive savings growth • Seasonal optimization • Continuous improvement benefits"
        
        p = insights_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.name = template_config['font_name']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = """
Savings Timeline talking points:
- Explain the ramp-up period for savings realization
- Discuss seasonal variations and how they're addressed
- Highlight the cumulative nature of benefits
- Address timeline expectations and milestones
- Emphasize conservative projections with upside potential
"""
    
    def _add_implementation_timeline_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add implementation timeline slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "Implementation Timeline & Milestones"
        self._format_title_textbox(title_frame, template_config)
        
        # Add Gantt chart
        if 'implementation_gantt' in self.generated_charts:
            chart_path = self.generated_charts['implementation_gantt']
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), Inches(12), Inches(4.5))
        
        # Key milestones
        milestones_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1))
        milestones_frame = milestones_box.text_frame
        
        milestones_frame.text = """Critical Milestones: Week 4 - System Ready | Week 12 - Go-Live | Week 24 - Full ROI Realization"""
        
        p = milestones_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.name = template_config['font_name']
        p.font.bold = True
        p.font.color.rgb = template_config['primary_color']
        p.alignment = PP_ALIGN.CENTER
        
        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = """
Implementation Timeline talking points:
- Walk through each phase and what it involves
- Highlight parallel activities to minimize timeline
- Discuss resource requirements from client side
- Address any concerns about business disruption
- Emphasize proven implementation methodology
"""
    
    def _add_risk_assessment_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add risk assessment slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "Risk Assessment & Mitigation"
        self._format_title_textbox(title_frame, template_config)
        
        # Left side - risk details
        left_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5))
        left_frame = left_textbox.text_frame
        
        left_frame.text = """Risk Mitigation Strategies:

Market Risk (Low)
• Conservative ROI projections
• Flexible solution architecture
• Proven track record in similar markets

Technology Risk (Low)
• Mature, tested technology platform
• Comprehensive backup and recovery
• 99.9% uptime SLA guarantee

Implementation Risk (Low-Medium)
• Experienced implementation team
• Phased rollout approach
• Dedicated client success manager

Financial Risk (Low)
• Performance-based milestone payments
• ROI guarantee with penalty clauses
• Transparent cost structure

Operational Risk (Low)
• Minimal business disruption
• Comprehensive training program
• 24/7 support during transition

Overall Risk Level: LOW
With appropriate mitigation strategies in place"""
        
        self._format_content_textbox(left_frame, template_config, font_size=Pt(11))
        
        # Right side - risk radar chart
        if 'risk_analysis' in self.generated_charts:
            chart_path = self.generated_charts['risk_analysis']
            slide.shapes.add_picture(chart_path, Inches(7), Inches(1.5), Inches(5.5), Inches(5))
        
        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = """
Risk Assessment talking points:
- Address each risk category honestly but confidently
- Provide specific examples of mitigation strategies
- Reference similar successful implementations
- Discuss contingency planning
- Emphasize low overall risk profile
- Be prepared to discuss any client-specific concerns
"""
    
    def _add_next_steps_slide(self, prs: Presentation, template_config: Dict) -> None:
        """Add next steps and recommendations slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Next Steps & Recommendations"
        self._format_title(title, template_config)
        
        # Content
        content = slide.placeholders[1]
        content.text = """Recommended Action Plan:

1. Decision & Approval
   • Review and approve ROI analysis
   • Secure internal stakeholder buy-in
   • Finalize budget allocation

2. Project Initiation (Week 1-2)
   • Execute service agreement
   • Conduct detailed discovery session
   • Establish project team and governance

3. Implementation Planning (Week 2-3)
   • Finalize technical requirements
   • Create detailed project timeline
   • Identify integration touchpoints

4. Go-Live Preparation (Week 3-12)
   • System configuration and testing
   • User training and change management
   • Data migration and validation

5. Success Monitoring
   • ROI tracking and reporting
   • Continuous optimization opportunities
   • Quarterly business reviews

Investment Decision Timeline: 2 weeks recommended
Implementation Start: Within 30 days of approval"""
        
        self._format_content(content, template_config, font_size=Pt(13))
        
        # Add contact information
        contact_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1))
        contact_frame = contact_box.text_frame
        
        contact_frame.text = f"Questions? Contact: {self.company_config['phone']} | {self.company_config['email']}"
        
        p = contact_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.name = template_config['font_name']
        p.font.bold = True
        p.font.color.rgb = template_config['accent_color']
        p.alignment = PP_ALIGN.CENTER
        
        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = """
Next Steps talking points:
- Create urgency around the decision timeline
- Emphasize support available throughout the process
- Address any remaining concerns or objections
- Discuss specific next steps and timeline
- Secure commitment to move forward
- Schedule follow-up meetings and technical discussions
"""
    
    def _format_title(self, title_shape, template_config: Dict) -> None:
        """Format title text"""
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = template_config['title_size']
        title_paragraph.font.name = template_config['font_name']
        title_paragraph.font.color.rgb = template_config['primary_color']
        title_paragraph.font.bold = True
    
    def _format_title_textbox(self, text_frame, template_config: Dict) -> None:
        """Format title textbox"""
        p = text_frame.paragraphs[0]
        p.font.size = template_config['title_size']
        p.font.name = template_config['font_name']
        p.font.color.rgb = template_config['primary_color']
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
    
    def _format_content(self, content_shape, template_config: Dict, font_size: Pt = None) -> None:
        """Format content text"""
        if font_size is None:
            font_size = template_config['body_size']
            
        for paragraph in content_shape.text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.name = template_config['font_name']
            paragraph.level = 0
    
    def _format_content_textbox(self, text_frame, template_config: Dict, font_size: Pt = None) -> None:
        """Format content textbox"""
        if font_size is None:
            font_size = template_config['body_size']
            
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.name = template_config['font_name']
    
    def get_available_templates(self) -> Dict[str, Dict]:
        """Get available presentation templates"""
        return {
            name: {
                'name': config['name'],
                'description': config['description']
            }
            for name, config in self.TEMPLATES.items()
        }
    
    def cleanup(self) -> None:
        """Clean up temporary files and directories"""
        try:
            import shutil
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
                logger.info(f"Cleaned up temporary directory: {self.temp_dir}")
        except Exception as e:
            logger.warning(f"Could not clean up temp directory {self.temp_dir}: {e}")
    
    def __del__(self):
        """Destructor to ensure cleanup"""
        self.cleanup()


if __name__ == '__main__':
    # Example usage
    sample_results = {
        'inputs': {
            'company_name': 'Acme E-commerce Ltd.',
            'annual_revenue': 2000000,
            'service_investment': 75000,
            'labor_costs': 8000,
            'shipping_costs': 5000,
            'error_costs': 2000,
            'inventory_costs': 3000
        },
        'roi_metrics': {
            'first_year_roi': 0.45,
            'annual_savings': 85000,
            'monthly_savings': 7083,
            'payback_period_months': 10.6
        },
        'financial_metrics': {
            'npv': 185000,
            'irr': 0.52
        },
        'savings': {
            'labor': {'annual': 34000, 'percentage': 0.6},
            'shipping': {'annual': 15000, 'percentage': 0.25},
            'errors': {'annual': 19000, 'percentage': 0.8},
            'inventory': {'annual': 11000, 'percentage': 0.3},
            'annual_total': 79000
        },
        'projections': {
            'year_1': {'roi_percentage': 0.45, 'cumulative_savings': 85000},
            'year_2': {'roi_percentage': 0.52, 'cumulative_savings': 178500},
            'year_3': {'roi_percentage': 0.58, 'cumulative_savings': 285000}
        }
    }
    
    # Create PowerPoint generator
    generator = PowerPointGenerator(sample_results)
    
    # Generate presentation
    for template in ['executive', 'sales', 'technical']:
        filepath = generator.generate_presentation(template)
        print(f"Generated {template} presentation: {filepath}")