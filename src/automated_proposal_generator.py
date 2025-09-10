#!/usr/bin/env python3
"""
Automated Proposal Generator for Chilean E-commerce Consulting
Generates professional proposals in 30 minutes vs 4 hours
"""

import json
import os
from datetime import datetime, timedelta
from typing import Dict, List, Optional
from dataclasses import dataclass
import pandas as pd
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_RIGHT, TA_JUSTIFY
from reportlab.pdfgen import canvas
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


@dataclass
class ProposalTemplate:
    """Structure for proposal templates"""
    name: str
    sections: List[str]
    style: str  # 'executive', 'technical', 'detailed'
    language: str = 'es'


class AutomatedProposalGenerator:
    """Generate professional consulting proposals automatically"""
    
    def __init__(self):
        self.templates = self._initialize_templates()
        self.case_studies = self._load_case_studies()
        self.service_packages = self._define_service_packages()
        self.proposal_data = {}
        
    def _initialize_templates(self) -> Dict[str, ProposalTemplate]:
        """Initialize proposal templates"""
        
        return {
            'executive': ProposalTemplate(
                name='Executive Proposal',
                sections=[
                    'executive_summary',
                    'business_case',
                    'roi_analysis',
                    'implementation_timeline',
                    'investment',
                    'next_steps'
                ],
                style='executive'
            ),
            'detailed': ProposalTemplate(
                name='Detailed Technical Proposal',
                sections=[
                    'executive_summary',
                    'current_state_analysis',
                    'pain_points',
                    'proposed_solution',
                    'technical_architecture',
                    'implementation_plan',
                    'roi_analysis',
                    'risk_mitigation',
                    'investment',
                    'timeline',
                    'team',
                    'next_steps'
                ],
                style='detailed'
            ),
            'quick': ProposalTemplate(
                name='Quick Proposal',
                sections=[
                    'executive_summary',
                    'key_benefits',
                    'solution_overview',
                    'investment',
                    'timeline',
                    'next_steps'
                ],
                style='quick'
            )
        }
    
    def _load_case_studies(self) -> List[Dict]:
        """Load relevant case studies"""
        
        return [
            {
                'client': 'Retail Fashion Chile',
                'industry': 'Retail',
                'challenge': 'Procesamiento manual de 2000+ órdenes mensuales',
                'solution': 'Automatización completa con integración Defontana-WooCommerce',
                'results': {
                    'time_saved': '75%',
                    'error_reduction': '92%',
                    'roi': '186%',
                    'payback': '5.2 meses'
                },
                'testimonial': 'La automatización transformó nuestras operaciones. Ahora procesamos el triple de órdenes con la mitad del equipo.'
            },
            {
                'client': 'Distribuidora Santiago',
                'industry': 'Wholesale',
                'challenge': 'Gestión de inventario en 5 canales de venta',
                'solution': 'Sistema centralizado con sincronización en tiempo real',
                'results': {
                    'stock_accuracy': '99.5%',
                    'overselling_eliminated': '100%',
                    'revenue_increase': '23%',
                    'roi': '210%'
                },
                'testimonial': 'Eliminamos completamente los quiebres de stock y la sobreventa. Nuestros clientes están más satisfechos que nunca.'
            },
            {
                'client': 'Servicios Profesionales Ltda',
                'industry': 'Services',
                'challenge': 'Facturación manual y seguimiento de proyectos',
                'solution': 'Automatización de facturación y portal de clientes',
                'results': {
                    'billing_time': '-90%',
                    'cash_flow': '+35%',
                    'client_satisfaction': '+40%',
                    'roi': '156%'
                },
                'testimonial': 'El portal de clientes mejoró dramáticamente nuestra relación con ellos y aceleró los pagos.'
            }
        ]
    
    def _define_service_packages(self) -> Dict[str, Dict]:
        """Define service packages with pricing tiers"""
        
        return {
            'starter': {
                'name': 'Paquete Inicio',
                'description': 'Ideal para empresas iniciando su transformación digital',
                'duration': '4-6 semanas',
                'price_clp': 8000000,
                'includes': [
                    'Diagnóstico operacional completo',
                    'Implementación de 2 automatizaciones clave',
                    'Integración básica de sistemas',
                    'Capacitación del equipo (8 horas)',
                    'Soporte por 30 días',
                    'Documentación de procesos'
                ],
                'best_for': 'Empresas con facturación < 30M CLP/mes'
            },
            'professional': {
                'name': 'Paquete Profesional',
                'description': 'Solución completa para optimización operacional',
                'duration': '8-10 semanas',
                'price_clp': 18000000,
                'includes': [
                    'Auditoría integral de operaciones',
                    'Automatización de todos los procesos críticos',
                    'Integración completa ERP-Ecommerce-Marketplace',
                    'Dashboard de KPIs en tiempo real',
                    'Capacitación intensiva (16 horas)',
                    'Soporte por 60 días',
                    'Optimización de fulfillment',
                    'Gestión de cambio organizacional'
                ],
                'best_for': 'Empresas con facturación 30-100M CLP/mes'
            },
            'enterprise': {
                'name': 'Paquete Enterprise',
                'description': 'Transformación digital completa',
                'duration': '12-16 semanas',
                'price_clp': 35000000,
                'includes': [
                    'Consultoría estratégica de transformación',
                    'Rediseño completo de operaciones',
                    'Automatización end-to-end',
                    'Integraciones avanzadas con IA',
                    'Sistema predictivo de inventario',
                    'Optimización multicanal',
                    'Dashboard ejecutivo personalizado',
                    'Capacitación continua (40 horas)',
                    'Soporte premium por 90 días',
                    'Garantía de resultados'
                ],
                'best_for': 'Empresas con facturación > 100M CLP/mes'
            },
            'custom': {
                'name': 'Paquete Personalizado',
                'description': 'Solución a medida según sus necesidades',
                'duration': 'A definir',
                'price_clp': 0,  # Custom pricing
                'includes': [
                    'Alcance personalizado según requerimientos',
                    'Combinación flexible de servicios',
                    'Timeline adaptado a su disponibilidad',
                    'Modelo de precios flexible',
                    'SLA personalizado'
                ],
                'best_for': 'Empresas con requerimientos específicos'
            }
        }
    
    def generate_proposal(self, 
                         client_data: Dict,
                         assessment_results: Dict,
                         roi_analysis: Dict,
                         template_type: str = 'executive',
                         package_type: str = 'professional') -> Dict:
        """
        Generate complete proposal based on inputs
        
        Args:
            client_data: Client information
            assessment_results: Results from rapid assessment
            roi_analysis: ROI calculation results
            template_type: Type of proposal template
            package_type: Service package to propose
        """
        
        template = self.templates[template_type]
        package = self.service_packages[package_type]
        
        # Build proposal data
        self.proposal_data = {
            'metadata': {
                'proposal_id': f"PROP-{datetime.now().strftime('%Y%m%d')}-{client_data.get('company_name', 'CLIENT')[:3].upper()}",
                'date': datetime.now().strftime('%d de %B de %Y'),
                'valid_until': (datetime.now() + timedelta(days=30)).strftime('%d de %B de %Y'),
                'consultant': 'Consultoría E-commerce Chile',
                'version': '1.0'
            },
            'client': client_data,
            'sections': {}
        }
        
        # Generate each section
        for section in template.sections:
            self.proposal_data['sections'][section] = self._generate_section(
                section, client_data, assessment_results, roi_analysis, package
            )
        
        # Add package details
        self.proposal_data['package'] = package
        
        # Add case study
        self.proposal_data['case_study'] = self._select_relevant_case_study(
            client_data.get('industry', 'Retail')
        )
        
        return self.proposal_data
    
    def _generate_section(self, section_name: str, client_data: Dict, 
                         assessment: Dict, roi: Dict, package: Dict) -> Dict:
        """Generate specific proposal section"""
        
        if section_name == 'executive_summary':
            return self._generate_executive_summary(client_data, assessment, roi)
        elif section_name == 'current_state_analysis':
            return self._generate_current_state(assessment)
        elif section_name == 'pain_points':
            return self._generate_pain_points(assessment)
        elif section_name == 'proposed_solution':
            return self._generate_solution(assessment, package)
        elif section_name == 'roi_analysis':
            return self._generate_roi_section(roi)
        elif section_name == 'implementation_timeline':
            return self._generate_timeline(package)
        elif section_name == 'investment':
            return self._generate_investment(package, roi)
        elif section_name == 'next_steps':
            return self._generate_next_steps()
        elif section_name == 'key_benefits':
            return self._generate_benefits(roi, assessment)
        elif section_name == 'technical_architecture':
            return self._generate_architecture(assessment)
        elif section_name == 'risk_mitigation':
            return self._generate_risk_mitigation()
        elif section_name == 'team':
            return self._generate_team()
        else:
            return {'title': section_name, 'content': 'Section under development'}
    
    def _generate_executive_summary(self, client: Dict, assessment: Dict, roi: Dict) -> Dict:
        """Generate executive summary section"""
        
        # Extract values with safe fallbacks
        improvements = roi.get('improvements', {})
        
        # Handle both nested and flat structures
        monthly_savings = improvements.get('total_monthly_savings_clp', 
                                          improvements.get('total_annual_savings_clp', 0) / 12 if improvements.get('total_annual_savings_clp') else 0)
        
        # Get efficiency metrics with fallbacks
        current_efficiency = roi.get('current_state', {}).get('operational_efficiency', 0.7)
        new_efficiency = improvements.get('new_operational_efficiency', current_efficiency + 0.15)
        efficiency_gain = (new_efficiency - current_efficiency) * 100
        
        # Get ROI metrics with fallbacks
        roi_percentage = improvements.get('roi_percentage_year_1', 
                                         improvements.get('roi_percentage', 100))
        payback_months = improvements.get('payback_months', 
                                         improvements.get('payback_period_months', 6))
        
        return {
            'title': 'Resumen Ejecutivo',
            'content': f"""
{client.get('company_name', 'Su empresa')} enfrenta desafíos operacionales que están limitando su crecimiento y rentabilidad en el competitivo mercado e-commerce chileno.

Nuestra evaluación identificó oportunidades para:
• Reducir costos operacionales en {monthly_savings/1000000:.1f}M CLP mensuales
• Mejorar la eficiencia operacional en {efficiency_gain:.0f}%
• Lograr un ROI de {roi_percentage:.0f}% en el primer año
• Recuperar la inversión en {payback_months:.1f} meses

Proponemos una solución integral que automatizará sus procesos críticos, eliminará las ineficiencias identificadas y posicionará a su empresa para un crecimiento sostenible.
            """,
            'highlights': [
                f"ROI proyectado: {roi_percentage:.0f}%",
                f"Ahorro anual: ${monthly_savings * 12/1000000:.1f}M CLP",
                f"Recuperación: {payback_months:.1f} meses"
            ]
        }
    
    def _generate_current_state(self, assessment: Dict) -> Dict:
        """Generate current state analysis section"""
        
        maturity = assessment.get('maturity_level', {})
        pain_points = assessment.get('pain_points', [])
        
        content = f"""
Nivel de Madurez Digital: {maturity.get('level', 'BÁSICO')}
{maturity.get('description', '')}

Evaluación por Área:
• Tecnología: {maturity.get('breakdown', {}).get('technology', 'N/A')}
• Operaciones: {maturity.get('breakdown', {}).get('operations', 'N/A')}
• Integración: {maturity.get('breakdown', {}).get('integration', 'N/A')}

Principales Desafíos Identificados:
"""
        
        for pain in pain_points[:3]:
            content += f"\n• {pain['issue']}: {pain['impact']}"
        
        return {
            'title': 'Análisis del Estado Actual',
            'content': content,
            'metrics': maturity.get('breakdown', {})
        }
    
    def _generate_pain_points(self, assessment: Dict) -> Dict:
        """Generate pain points section"""
        
        pain_points = assessment.get('pain_points', [])
        
        points = []
        total_cost = 0
        
        for pain in pain_points:
            points.append({
                'issue': pain['issue'],
                'severity': pain['severity'],
                'impact': pain['impact'],
                'monthly_cost': pain.get('cost_impact_clp', 0)
            })
            total_cost += pain.get('cost_impact_clp', 0)
        
        return {
            'title': 'Problemas Identificados y Su Impacto',
            'points': points,
            'total_monthly_cost': total_cost,
            'annual_cost': total_cost * 12,
            'summary': f"Estos problemas están costando aproximadamente ${total_cost/1000000:.1f}M CLP mensuales"
        }
    
    def _generate_solution(self, assessment: Dict, package: Dict) -> Dict:
        """Generate proposed solution section"""
        
        opportunities = assessment.get('opportunities', [])
        recommendations = assessment.get('recommendations', [])
        
        solution_components = []
        
        for rec in recommendations[:5]:
            solution_components.append({
                'component': rec['title'],
                'description': rec['description'],
                'impact': rec['expected_impact'],
                'timeline': rec.get('implementation_time', 'Por definir')
            })
        
        return {
            'title': 'Solución Propuesta',
            'overview': f"""
Implementaremos una solución integral que aborda sus desafíos operacionales mediante:

1. **Automatización de Procesos**: Eliminación de tareas manuales repetitivas
2. **Integración de Sistemas**: Conexión seamless entre sus plataformas
3. **Optimización Operacional**: Mejora de eficiencia en todos los procesos
4. **Capacitación del Equipo**: Empoderamiento de su personal con nuevas herramientas
5. **Monitoreo Continuo**: Dashboard de KPIs para toma de decisiones
            """,
            'components': solution_components,
            'package_name': package['name'],
            'package_includes': package['includes']
        }
    
    def _generate_roi_section(self, roi: Dict) -> Dict:
        """Generate ROI analysis section"""
        
        scenarios = roi.get('scenarios', {}).get('scenarios', {})
        monte_carlo = roi.get('scenarios', {}).get('monte_carlo', {})
        
        return {
            'title': 'Análisis de Retorno de Inversión',
            'summary': roi.get('executive_summary', {}),
            'scenarios': {
                'conservative': {
                    'name': 'Escenario Conservador',
                    'roi': scenarios.get('pessimistic', {}).get('roi_percentage', 0),
                    'savings': scenarios.get('pessimistic', {}).get('annual_savings_clp', 0),
                    'probability': scenarios.get('pessimistic', {}).get('probability', 0)
                },
                'expected': {
                    'name': 'Escenario Esperado',
                    'roi': scenarios.get('realistic', {}).get('roi_percentage', 0),
                    'savings': scenarios.get('realistic', {}).get('annual_savings_clp', 0),
                    'probability': scenarios.get('realistic', {}).get('probability', 0)
                },
                'optimistic': {
                    'name': 'Escenario Optimista',
                    'roi': scenarios.get('optimistic', {}).get('roi_percentage', 0),
                    'savings': scenarios.get('optimistic', {}).get('annual_savings_clp', 0),
                    'probability': scenarios.get('optimistic', {}).get('probability', 0)
                }
            },
            'confidence': monte_carlo.get('probability_positive_roi', 95),
            'three_year_projection': roi.get('three_year_projection', {})
        }
    
    def _generate_timeline(self, package: Dict) -> Dict:
        """Generate implementation timeline"""
        
        duration = package.get('duration', '8-10 semanas')
        
        # Parse duration to get weeks
        import re
        weeks_match = re.search(r'(\d+)-(\d+)', duration)
        if weeks_match:
            min_weeks = int(weeks_match.group(1))
            max_weeks = int(weeks_match.group(2))
        else:
            min_weeks, max_weeks = 8, 10
        
        phases = [
            {
                'phase': 'Fase 1: Descubrimiento y Diseño',
                'duration': f'{min_weeks//4} semanas',
                'activities': [
                    'Análisis detallado de procesos actuales',
                    'Mapeo de sistemas y integraciones',
                    'Diseño de solución y arquitectura',
                    'Definición de KPIs y métricas'
                ]
            },
            {
                'phase': 'Fase 2: Implementación Core',
                'duration': f'{min_weeks//2} semanas',
                'activities': [
                    'Configuración de automatizaciones',
                    'Desarrollo de integraciones',
                    'Implementación de validaciones',
                    'Setup de monitoreo'
                ]
            },
            {
                'phase': 'Fase 3: Testing y Optimización',
                'duration': f'{min_weeks//4} semanas',
                'activities': [
                    'Testing exhaustivo de procesos',
                    'Ajustes y optimizaciones',
                    'Capacitación del equipo',
                    'Documentación final'
                ]
            },
            {
                'phase': 'Fase 4: Go-Live y Estabilización',
                'duration': '1 semana',
                'activities': [
                    'Migración a producción',
                    'Monitoreo intensivo',
                    'Soporte on-site',
                    'Handover al equipo'
                ]
            }
        ]
        
        return {
            'title': 'Timeline de Implementación',
            'total_duration': duration,
            'start_date': 'A confirmar',
            'phases': phases,
            'milestones': [
                {'week': 2, 'milestone': 'Diseño aprobado'},
                {'week': min_weeks//2, 'milestone': 'Automatizaciones operativas'},
                {'week': min_weeks-1, 'milestone': 'Testing completado'},
                {'week': max_weeks, 'milestone': 'Go-live exitoso'}
            ]
        }
    
    def _generate_investment(self, package: Dict, roi: Dict) -> Dict:
        """Generate investment section"""
        
        price = package['price_clp']
        
        # Calculate payment options
        payment_options = [
            {
                'option': 'Pago Contado',
                'description': '100% al inicio del proyecto',
                'amount': price,
                'discount': price * 0.05,  # 5% discount
                'total': price * 0.95
            },
            {
                'option': 'Pago 50/50',
                'description': '50% inicial, 50% al go-live',
                'payments': [
                    {'milestone': 'Firma de contrato', 'amount': price * 0.5},
                    {'milestone': 'Go-live', 'amount': price * 0.5}
                ],
                'total': price
            },
            {
                'option': 'Pago en 3 Cuotas',
                'description': '40% inicial, 30% intermedio, 30% final',
                'payments': [
                    {'milestone': 'Firma de contrato', 'amount': price * 0.4},
                    {'milestone': 'Fin Fase 2', 'amount': price * 0.3},
                    {'milestone': 'Go-live', 'amount': price * 0.3}
                ],
                'total': price
            }
        ]
        
        # Calculate financing metrics
        monthly_savings = roi.get('improvements', {}).get('total_monthly_savings_clp', 0)
        months_to_pay = price / monthly_savings if monthly_savings > 0 else float('inf')
        
        return {
            'title': 'Inversión',
            'package_name': package['name'],
            'base_price': price,
            'payment_options': payment_options,
            'included_services': package['includes'],
            'not_included': [
                'Licencias de software de terceros',
                'Hardware adicional si se requiere',
                'Modificaciones fuera del alcance acordado',
                'Soporte posterior al período incluido'
            ],
            'roi_context': {
                'monthly_savings': monthly_savings,
                'months_to_recover': roi.get('improvements', {}).get('payback_months', 0),
                'self_funding_months': months_to_pay
            },
            'validity': '30 días desde la fecha de esta propuesta'
        }
    
    def _generate_next_steps(self) -> Dict:
        """Generate next steps section"""
        
        return {
            'title': 'Próximos Pasos',
            'steps': [
                {
                    'step': 1,
                    'action': 'Revisión de Propuesta',
                    'description': 'Revisaremos juntos esta propuesta para aclarar cualquier duda',
                    'timeline': 'Esta semana'
                },
                {
                    'step': 2,
                    'action': 'Ajustes y Aprobación',
                    'description': 'Ajustaremos la propuesta según su feedback y procederemos con la aprobación',
                    'timeline': 'Próximos 3-5 días'
                },
                {
                    'step': 3,
                    'action': 'Firma de Contrato',
                    'description': 'Formalizaremos el acuerdo con términos y condiciones claros',
                    'timeline': 'Próxima semana'
                },
                {
                    'step': 4,
                    'action': 'Kick-off del Proyecto',
                    'description': 'Reunión de inicio con todos los stakeholders',
                    'timeline': '3 días después de la firma'
                },
                {
                    'step': 5,
                    'action': 'Inicio de Implementación',
                    'description': 'Comenzamos con la fase de descubrimiento y análisis',
                    'timeline': 'Inmediatamente después del kick-off'
                }
            ],
            'call_to_action': """
¿Listo para transformar sus operaciones?

Contáctenos:
📧 Email: consulting@ecommerce.cl
📱 Teléfono: +56 9 XXXX XXXX
🌐 Web: www.ecommerceconsulting.cl

Agendemos una reunión esta semana para discutir cómo podemos ayudarle a alcanzar sus objetivos.
            """
        }
    
    def _generate_benefits(self, roi: Dict, assessment: Dict) -> Dict:
        """Generate key benefits section"""
        
        return {
            'title': 'Beneficios Clave',
            'benefits': [
                {
                    'category': 'Financieros',
                    'items': [
                        f"Ahorro de ${roi['improvements']['total_annual_savings_clp']/1000000:.1f}M CLP anuales",
                        f"ROI de {roi['improvements']['roi_percentage_year_1']:.0f}% en año 1",
                        f"Recuperación de inversión en {roi['improvements']['payback_months']:.1f} meses"
                    ]
                },
                {
                    'category': 'Operacionales',
                    'items': [
                        'Reducción del 70% en tiempo de procesamiento',
                        'Eliminación del 90% de errores manuales',
                        'Automatización de procesos repetitivos'
                    ]
                },
                {
                    'category': 'Estratégicos',
                    'items': [
                        'Escalabilidad para crecimiento futuro',
                        'Mejor toma de decisiones con data en tiempo real',
                        'Ventaja competitiva en el mercado'
                    ]
                }
            ]
        }
    
    def _select_relevant_case_study(self, industry: str) -> Dict:
        """Select most relevant case study"""
        
        # Find matching industry
        for case in self.case_studies:
            if case['industry'].lower() == industry.lower():
                return case
        
        # Return first if no match
        return self.case_studies[0] if self.case_studies else {}
    
    def _generate_architecture(self, assessment: Dict) -> Dict:
        """Generate technical architecture section"""
        
        return {
            'title': 'Arquitectura Técnica',
            'components': [
                {
                    'layer': 'Capa de Presentación',
                    'technologies': ['Dashboard Web', 'APIs REST', 'Webhooks']
                },
                {
                    'layer': 'Capa de Integración',
                    'technologies': ['API Gateway', 'Message Queue', 'ETL Pipelines']
                },
                {
                    'layer': 'Capa de Procesamiento',
                    'technologies': ['Automation Engine', 'Business Rules', 'Validation Layer']
                },
                {
                    'layer': 'Capa de Datos',
                    'technologies': ['PostgreSQL', 'Redis Cache', 'Data Warehouse']
                }
            ],
            'integrations': [
                'Defontana/Bsale (ERP)',
                'WooCommerce/Shopify (E-commerce)',
                'Transbank/Webpay (Payments)',
                'Chilexpress/Starken (Shipping)',
                'MercadoLibre/Falabella (Marketplaces)'
            ]
        }
    
    def _generate_risk_mitigation(self) -> Dict:
        """Generate risk mitigation section"""
        
        return {
            'title': 'Mitigación de Riesgos',
            'risks': [
                {
                    'risk': 'Resistencia al cambio del equipo',
                    'probability': 'Media',
                    'impact': 'Alto',
                    'mitigation': 'Plan de gestión del cambio con capacitación intensiva y soporte continuo'
                },
                {
                    'risk': 'Complejidad técnica de integraciones',
                    'probability': 'Baja',
                    'impact': 'Medio',
                    'mitigation': 'Equipo experto en integraciones chilenas, ambiente de testing robusto'
                },
                {
                    'risk': 'Interrupción de operaciones',
                    'probability': 'Muy Baja',
                    'impact': 'Alto',
                    'mitigation': 'Implementación gradual, rollback plan, trabajo fuera de horario peak'
                },
                {
                    'risk': 'Sobrecostos',
                    'probability': 'Baja',
                    'impact': 'Medio',
                    'mitigation': 'Precio fijo cerrado, alcance bien definido, change management process'
                }
            ]
        }
    
    def _generate_team(self) -> Dict:
        """Generate team section"""
        
        return {
            'title': 'Equipo del Proyecto',
            'members': [
                {
                    'role': 'Project Manager',
                    'name': 'Por asignar',
                    'experience': '10+ años en transformación digital',
                    'certifications': ['PMP', 'Agile']
                },
                {
                    'role': 'Arquitecto de Soluciones',
                    'name': 'Por asignar',
                    'experience': '8+ años en integraciones e-commerce',
                    'expertise': ['API Design', 'System Integration']
                },
                {
                    'role': 'Desarrollador Senior',
                    'name': 'Por asignar',
                    'experience': '5+ años en automatización',
                    'technologies': ['Python', 'Node.js', 'SQL']
                },
                {
                    'role': 'Consultor de Procesos',
                    'name': 'Por asignar',
                    'experience': '7+ años en optimización operacional',
                    'specialization': 'E-commerce operations'
                }
            ]
        }
    
    def export_to_pdf(self, filename: str = 'proposal.pdf'):
        """Export proposal to PDF"""
        
        doc = SimpleDocTemplate(filename, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1e40af'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#1e40af'),
            spaceBefore=20,
            spaceAfter=12
        )
        
        # Cover page
        story.append(Paragraph('PROPUESTA DE CONSULTORÍA', title_style))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"Para: {self.proposal_data['client'].get('company_name', 'Cliente')}", styles['Normal']))
        story.append(Paragraph(f"Fecha: {self.proposal_data['metadata']['date']}", styles['Normal']))
        story.append(Paragraph(f"Propuesta ID: {self.proposal_data['metadata']['proposal_id']}", styles['Normal']))
        story.append(PageBreak())
        
        # Add sections
        for section_key, section_data in self.proposal_data['sections'].items():
            story.append(Paragraph(section_data.get('title', ''), heading_style))
            
            content = section_data.get('content', '')
            if content:
                # Split content into paragraphs
                for para in content.split('\n\n'):
                    if para.strip():
                        story.append(Paragraph(para.strip(), styles['Normal']))
                        story.append(Spacer(1, 0.1*inch))
            
            # Add any additional structured data
            if 'highlights' in section_data:
                for highlight in section_data['highlights']:
                    story.append(Paragraph(f"• {highlight}", styles['Normal']))
                story.append(Spacer(1, 0.2*inch))
            
            if 'components' in section_data:
                for comp in section_data['components']:
                    story.append(Paragraph(f"<b>{comp.get('component', '')}</b>", styles['Normal']))
                    story.append(Paragraph(comp.get('description', ''), styles['Normal']))
                    story.append(Spacer(1, 0.1*inch))
        
        # Build PDF
        doc.build(story)
        
    def export_to_powerpoint(self, filename: str = 'proposal.pptx'):
        """Export proposal to PowerPoint"""
        
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = "Propuesta de Consultoría E-commerce"
        subtitle.text = f"{self.proposal_data['client'].get('company_name', 'Cliente')}\n{self.proposal_data['metadata']['date']}"
        
        # Executive Summary slide
        if 'executive_summary' in self.proposal_data['sections']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Resumen Ejecutivo"
            
            content = slide.placeholders[1]
            summary = self.proposal_data['sections']['executive_summary']
            
            # Add highlights as bullet points
            if 'highlights' in summary:
                bullet_text = '\n'.join([f"• {h}" for h in summary['highlights']])
                content.text = bullet_text
        
        # ROI slide
        if 'roi_analysis' in self.proposal_data['sections']:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
            title = slide.shapes.title
            title.text = "Análisis de ROI"
            
            # Add ROI metrics as text boxes
            roi_data = self.proposal_data['sections']['roi_analysis']
            scenarios = roi_data.get('scenarios', {})
            
            # Add expected scenario
            if 'expected' in scenarios:
                expected = scenarios['expected']
                textbox = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(8), Inches(1)
                )
                tf = textbox.text_frame
                tf.text = f"ROI Esperado: {expected['roi']:.0f}%"
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.bold = True
        
        # Investment slide
        if 'investment' in self.proposal_data['sections']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Inversión"
            
            content = slide.placeholders[1]
            investment = self.proposal_data['sections']['investment']
            
            text = f"Paquete: {investment['package_name']}\n"
            text += f"Inversión: ${investment['base_price']:,.0f} CLP\n\n"
            text += "Incluye:\n"
            for service in investment['included_services'][:5]:
                text += f"• {service}\n"
            
            content.text = text
        
        # Next steps slide
        if 'next_steps' in self.proposal_data['sections']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Próximos Pasos"
            
            content = slide.placeholders[1]
            steps = self.proposal_data['sections']['next_steps']['steps']
            
            text = ""
            for step in steps[:5]:
                text += f"{step['step']}. {step['action']}\n"
            
            content.text = text
        
        # Save presentation
        prs.save(filename)
    
    def generate_one_pager(self) -> str:
        """Generate executive one-pager summary"""
        
        if not self.proposal_data:
            return "No proposal data available"
        
        roi = self.proposal_data['sections'].get('roi_analysis', {})
        investment = self.proposal_data['sections'].get('investment', {})
        
        one_pager = f"""
{'='*60}
RESUMEN EJECUTIVO - {self.proposal_data['client'].get('company_name', 'CLIENTE')}
{'='*60}

OPORTUNIDAD
-----------
Potencial de ahorro: ${roi.get('summary', {}).get('annual_savings_clp', 0)/1000000:.1f}M CLP/año
ROI proyectado: {roi.get('summary', {}).get('headline_roi', 0):.0f}%
Recuperación: {roi.get('summary', {}).get('payback_period_months', 0):.1f} meses

SOLUCIÓN
--------
{investment.get('package_name', 'Paquete Profesional')}
Duración: {self.proposal_data['package'].get('duration', '8-10 semanas')}
Inversión: ${investment.get('base_price', 0)/1000000:.1f}M CLP

BENEFICIOS CLAVE
---------------
✓ Automatización completa de procesos
✓ Reducción 70% en tiempo operativo
✓ Eliminación 90% de errores
✓ Dashboard KPIs tiempo real
✓ ROI garantizado

PRÓXIMO PASO
-----------
Agendar reunión esta semana para revisar propuesta detallada.

Contacto: consulting@ecommerce.cl | +56 9 XXXX XXXX
{'='*60}
"""
        
        return one_pager


# Example usage
if __name__ == "__main__":
    # Sample client data
    client_data = {
        'company_name': 'Tienda Online Chile SpA',
        'contact_name': 'Juan Pérez',
        'email': 'juan@tiendaonline.cl',
        'phone': '+56 9 8765 4321',
        'industry': 'Retail',
        'website': 'www.tiendaonline.cl'
    }
    
    # Sample assessment results (from rapid assessment tool)
    assessment_results = {
        'maturity_level': {
            'level': 'BÁSICO',
            'score': 4.5,
            'description': 'Procesos mayormente manuales con automatización limitada',
            'breakdown': {
                'technology': '3.5/10',
                'operations': '4.0/10',
                'integration': '5.5/10'
            }
        },
        'pain_points': [
            {
                'issue': 'Procesamiento Manual de Órdenes',
                'severity': 'ALTA',
                'impact': '20 horas semanales en tareas manuales',
                'cost_impact_clp': 3500000
            },
            {
                'issue': 'Alta Tasa de Errores',
                'severity': 'CRÍTICA',
                'impact': '8% de órdenes con errores',
                'cost_impact_clp': 2000000
            }
        ],
        'opportunities': [
            {
                'area': 'Automatización de Procesos',
                'monthly_savings_clp': 5000000,
                'implementation_effort': 'MEDIO',
                'time_to_value': '2-3 semanas'
            }
        ],
        'recommendations': [
            {
                'title': 'Implementar Integración ERP-Ecommerce',
                'description': 'Conectar Defontana con WooCommerce',
                'expected_impact': 'Ahorro de 20 horas semanales',
                'implementation_time': '3-4 semanas'
            }
        ]
    }
    
    # Sample ROI analysis (from enhanced ROI calculator)
    roi_analysis = {
        'current_state': {
            'operational_efficiency': 0.65
        },
        'improvements': {
            'total_monthly_savings_clp': 8500000,
            'total_annual_savings_clp': 102000000,
            'roi_percentage_year_1': 186,
            'payback_months': 5.2,
            'new_operational_efficiency': 0.85
        },
        'scenarios': {
            'scenarios': {
                'pessimistic': {
                    'roi_percentage': 120,
                    'annual_savings_clp': 70000000,
                    'probability': 25
                },
                'realistic': {
                    'roi_percentage': 186,
                    'annual_savings_clp': 102000000,
                    'probability': 60
                },
                'optimistic': {
                    'roi_percentage': 250,
                    'annual_savings_clp': 140000000,
                    'probability': 15
                }
            },
            'monte_carlo': {
                'probability_positive_roi': 98.5
            }
        },
        'executive_summary': {
            'headline_roi': 186,
            'payback_period_months': 5.2,
            'annual_savings_clp': 102000000
        },
        'three_year_projection': {
            'year_1': {'roi_percentage': 186},
            'year_2': {'roi_percentage': 420},
            'year_3': {'roi_percentage': 680}
        }
    }
    
    # Generate proposal
    generator = AutomatedProposalGenerator()
    proposal = generator.generate_proposal(
        client_data=client_data,
        assessment_results=assessment_results,
        roi_analysis=roi_analysis,
        template_type='executive',
        package_type='professional'
    )
    
    # Generate outputs
    print(generator.generate_one_pager())
    
    # Export to files
    generator.export_to_pdf('proposal.pdf')
    generator.export_to_powerpoint('proposal.pptx')
    
    # Save JSON
    with open('proposal_data.json', 'w', encoding='utf-8') as f:
        json.dump(proposal, f, ensure_ascii=False, indent=2, default=str)
    
    print("\n✅ Proposal generated successfully!")
    print("Files created: proposal.pdf, proposal.pptx, proposal_data.json")